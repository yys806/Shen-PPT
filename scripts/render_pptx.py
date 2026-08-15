#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
render_pptx.py — PPT 专属 skill 渲染引擎（v1.0）
读 deck-spec.json（内容单源）+ style-spec.json（风格 token）+ skeleton json（版式骨架）
→ 用 python-pptx 参数化渲染原生可编辑 .pptx

用法:
  python render_pptx.py -s deck-spec.json -o output.pptx [--style <name>] [--print-spec]
依赖: python-pptx (pip install python-pptx); PIL 可选(图片比例 fit 更准)
"""
import argparse
import json
import math
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    sys.exit("缺少依赖: pip install python-pptx")

SKILL_ROOT = Path(__file__).resolve().parent.parent
EMU_PER_PX = 9525  # 1px @96dpi = 9525 EMU
PX_TO_PT = 0.75     # 1px = 0.75pt (96dpi)
WARN = []


# ---------- 工具 ----------

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def split_ascii_runs(text, font_cn, font_en):
    """把文本按 ASCII/非 ASCII 分段，返回 (segment, use_en_font) 列表，保证中西文字体分工。"""
    runs = []
    for seg in re.split(r"([\x00-\x7F]+)", text):
        if not seg:
            continue
        runs.append((seg, seg.isascii()))
    return runs


def estimate_text_height(text_list, font_size_pt, region_w_px):
    """粗略估算多段文本总高度(pt)。font_size_pt 为字号，中文字宽≈字号。"""
    cpl = max(4, int(region_w_px / (font_size_pt * 1.333)))
    total_lines = 0
    for para in text_list:
        if not para:
            total_lines += 1
            continue
        lines = math.ceil(len(para) / cpl)
        total_lines += max(1, lines)
    return total_lines * font_size_pt * 1.25


def warn(msg):
    WARN.append(msg)
    print(f"  ⚠️ {msg}")


# ---------- 配置加载 ----------

def load_json(path):
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def resolve_style(spec, style_override):
    index = load_json(SKILL_ROOT / "styles" / "index.json")
    scenario = spec.get("scenario")
    style_name = style_override or spec.get("style")

    if not style_name:
        # 从 scenario 反查
        sc = index["scenarios"].get(scenario)
        if not sc:
            sys.exit(f"✗ scenario '{scenario}' 未在 styles/index.json 注册")
        style_name = sc["style"]

    entry = index["styles"].get(style_name)
    if not entry:
        sys.exit(f"✗ style '{style_name}' 未在 styles/index.json 注册")

    spec_path = SKILL_ROOT / "styles" / entry["path"]
    style = load_json(spec_path)

    # 容器风格（如 course-showcase）不能直接渲染
    if "children" in style:
        sys.exit(f"✗ '{style_name}' 是风格容器，请指定子风格: {', '.join(style['children'])}")

    # 锁定校验：scenario 锁定风格时不允许换
    sc_entry = index["scenarios"].get(scenario)
    if sc_entry and sc_entry.get("locked") and style_name != sc_entry["style"]:
        sys.exit(f"✗ scenario '{scenario}' 锁定风格 '{sc_entry['style']}'，不能改用 '{style_name}'")
    return style, style_name


def resolve_skeleton(style):
    sk_name = style.get("skeleton", "defense")
    path = SKILL_ROOT / "layouts" / f"skeleton-{sk_name}.json"
    if not path.exists():
        sys.exit(f"✗ 骨架文件不存在: {path}")
    sk = load_json(path)
    # 风格级骨架覆盖（skeleton_overrides）：bg / regions 按 slide 类型合并
    ov = style.get("skeleton_overrides") or {}
    for stype, sdict in ov.items():
        if stype not in sk.get("slides", {}):
            continue
        target = sk["slides"][stype]
        if "bg" in sdict:
            target["bg"] = {**target.get("bg", {}), **sdict["bg"]}
        if "regions" in sdict:
            target["regions"] = {**target.get("regions", {}), **sdict["regions"]}
    return sk


def validate_spec(spec):
    """内容单源校验。"""
    errs = []
    if not spec.get("deckTitle"):
        errs.append("缺少 deckTitle")
    scenario = spec.get("scenario")
    if scenario not in ("s1-pku", "s1-tongji", "s2-course", "s3-thesis", "s4-html"):
        errs.append(f"scenario '{scenario}' 非法（s1-pku/s1-tongji/s2-course/s3-thesis/s4-html）")
    sections = {s["id"] for s in spec.get("sections", [])}
    if scenario in ("s2-course", "s3-thesis") and not sections:
        errs.append("S2/S3 必须定义 sections[]")
    for i, sl in enumerate(spec.get("slides", [])):
        t = sl.get("type")
        if t not in ("cover", "toc", "section", "content", "ending"):
            errs.append(f"slide[{i}] type '{t}' 非法")
        if t in ("section", "content") and scenario in ("s2-course", "s3-thesis"):
            if not sl.get("section") or sl["section"] not in sections:
                errs.append(f"slide[{i}] ({t}) 缺少合法 section（S2/S3 硬规则）")
        if t == "content" and sl.get("layout") not in (None, "text", "bullets", "image-right", "image-left", "two-col", "table", "stats"):
            errs.append(f"slide[{i}] layout '{sl.get('layout')}' 非法")
    if errs:
        sys.exit("✗ deck-spec 校验失败:\n  - " + "\n  - ".join(errs))


# ---------- 渲染 ----------

class Renderer:
    def __init__(self, spec, style, skeleton):
        self.spec = spec
        self.style = style
        self.sk = skeleton
        self.pal = style["palette"]
        self.fonts = style["fonts"]
        self.sizes = style["sizes"]
        self.canvas = skeleton["canvas"]
        self.prs = Presentation()
        self.prs.slide_width = Emu(self.canvas["w"] * EMU_PER_PX)
        self.prs.slide_height = Emu(self.canvas["h"] * EMU_PER_PX)
        self.blank = self.prs.slide_layouts[6]
        self.page_no = 0  # 显示页码（cover 不计）
        self.sections = {s["id"]: s for s in spec.get("sections", [])}
        self.spec_dir = Path(spec.get("_spec_path", ".")).parent

    def size(self, token, default=18):
        return self.sizes.get(token, default)

    def font(self, token, default=None):
        return self.fonts.get(token, self.fonts.get("cn", "微软雅黑"))

    # -- 基础元素 --

    def add_rect(self, slide, x, y, w, h, fill_hex, line_hex=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x * EMU_PER_PX), Emu(y * EMU_PER_PX),
                                     Emu(w * EMU_PER_PX), Emu(h * EMU_PER_PX))
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        if line_hex:
            shp.line.color.rgb = hex_to_rgb(line_hex)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def add_text(self, slide, region, text, size_token=None, color_token=None, bold=None,
                 align=None, anchor=MSO_ANCHOR.TOP, font_override=None, sup=True):
        x, y, w, h = region["x"], region["y"], region["w"], region["h"]
        size = self.size(size_token or region.get("size"), 18)
        color = hex_to_rgb(self.pal[color_token or region.get("color") or
                                    ("accent" if region.get("accent") else
                                     ("muted" if region.get("muted") else "text"))])
        bold = region.get("bold", False) if bold is None else bold
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        align = align_map.get(align or region.get("align", "left"), PP_ALIGN.LEFT)

        box = slide.shapes.add_textbox(Emu(x * EMU_PER_PX), Emu(y * EMU_PER_PX),
                                       Emu(w * EMU_PER_PX), Emu(h * EMU_PER_PX))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        paragraphs = text if isinstance(text, list) else [text]
        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = 1.15
            cn_font = font_override or (self.fonts.get(region.get("font"), region.get("font")) if region.get("font") else None) or (self.font("cn_title") if size_token == "title" else self.font("cn"))
            for seg, is_en in split_ascii_runs(para, self.font("cn"), self.font("en")):
                self.add_smart_run(p, seg, size, color, bold, cn_font, self.font("en"), sup=sup)
        return box

    def parse_bullet_items(self, body, bullets_cfg):
        """把 body 列表解析为富文本项。支持 markdown 式层级：'# '一级 / '## '二级 / '### '三级。
        bullets_cfg: style-spec 的 bullets 配置 {l0/l1/l2/l3: {char, bullet_font, mark_color,
        size, color, bold, lvl, marL, lum_mod, lum_off}} —— 参数直接来自单位模板的原生段落定义。
        渲染为 PowerPoint 原生项目符号（Wingdings + 大纲级别 + 悬挂缩进），编辑时回车自动延续层级。"""
        out = []
        for it in body:
            s = str(it)
            m = re.match(r"^(#{1,3})\s*(.*)$", s)
            if m:
                level = len(m.group(1))
                text = m.group(2)
            else:
                level = 0
                text = s
            cfg = bullets_cfg.get(f"l{level}") or bullets_cfg.get("l0", {})
            item = dict(cfg)
            item["text"] = text
            out.append(item)
        return out

    def set_native_bullet(self, p, cfg):
        """写原生 PowerPoint 段落格式：大纲级别 lvl + 悬挂缩进 marL/indent +
        Wingdings 项目符号（buClr/buSzTx/buFont/buChar）。回车后自动延续层级。"""
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lvl = cfg.get("lvl")
        marL = cfg.get("marL")
        if lvl is not None:
            pPr.set("lvl", str(lvl))
        if marL:
            pPr.set("marL", str(marL))
            pPr.set("indent", str(-(cfg.get("indent", 228600))))
        if cfg.get("mark_color"):
            buClr = pPr.makeelement(qn("a:buClr"), {})
            srgb = pPr.makeelement(qn("a:srgbClr"), {"val": cfg["mark_color"]})
            if cfg.get("lum_mod") is not None:
                srgb.append(pPr.makeelement(qn("a:lumMod"), {"val": str(cfg["lum_mod"])}))
                srgb.append(pPr.makeelement(qn("a:lumOff"), {"val": str(cfg["lum_off"])}))
            buClr.append(srgb)
            pPr.append(buClr)
            pPr.append(pPr.makeelement(qn("a:buSzTx"), {}))
            pPr.append(pPr.makeelement(qn("a:buFont"),
                                       {"typeface": cfg.get("bullet_font", "Wingdings"),
                                        "pitchFamily": "2", "charset": "2"}))
            if cfg.get("char"):
                pPr.append(pPr.makeelement(qn("a:buChar"), {"char": cfg["char"]}))

    def build_runs(self, text, cfg):
        """把条目文本拆成 run 级片段 [(seg, color_token, bold), ...]。
        规则：
        - 层级2/3（lvl 1/2）：冒号前短语（≤12字）= 总结性小标题 → 蓝色加粗；
          冒号后 = 说明文字 → 黑色不加粗
        - **关键词** → 红色加粗（任意位置，优先于冒号规则）
        - 无冒号的三级条目：整条蓝色加粗（模板）；正文/一级：默认格式"""
        runs = []
        split = None
        if cfg.get("lvl") in (1, 2):
            m = re.match(r"^([^：:]{1,24}?)([：:])(.*)$", text)
            if m:
                split = (m.group(1), m.group(2), m.group(3))
        if split:
            pre, colon, post = split
            # 冒号前：** 红，否则蓝加粗
            for seg in re.split(r"(\*\*[^*]+\*\*)", pre):
                if not seg:
                    continue
                if seg.startswith("**") and seg.endswith("**"):
                    runs.append((seg[2:-2], "red", True))
                else:
                    runs.append((seg, "accent", True))
            # 冒号后：** 红，否则黑不加粗
            for seg in re.split(r"(\*\*[^*]+\*\*)", colon + post):
                if not seg:
                    continue
                if seg.startswith("**") and seg.endswith("**"):
                    runs.append((seg[2:-2], "red", True))
                else:
                    runs.append((seg, "text", False))
        else:
            for seg in re.split(r"(\*\*[^*]+\*\*)", text):
                if not seg:
                    continue
                if seg.startswith("**") and seg.endswith("**"):
                    runs.append((seg[2:-2], "red", True))
                else:
                    runs.append((seg, cfg["color"], cfg.get("bold", False)))
        return runs

    def add_rich_paras(self, slide, region, items, align=None):
        """渲染富文本段落列表（原生项目符号 + run 级强调）。items=[{text,size,color,bold,char,mark_color,lvl,marL,...}]"""
        x, y, w, h = region["x"], region["y"], region["w"], region["h"]
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        align = align_map.get(align or region.get("align", "left"), PP_ALIGN.LEFT)
        box = slide.shapes.add_textbox(Emu(x * EMU_PER_PX), Emu(y * EMU_PER_PX),
                                       Emu(w * EMU_PER_PX), Emu(h * EMU_PER_PX))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = 1.4  # 行距（页面更饱满）
            self.set_native_bullet(p, it)
            size = self.size(it["size"], 18)
            for seg, color_tok, bold in self.build_runs(it["text"], it):
                self.add_smart_run(p, seg, size, hex_to_rgb(self.pal[color_tok]), bold,
                                   self.font("cn"), self.font("en"))
        return box

    def render_logo(self, slide, region):
        """渲染右上角 LOGO 组。region={"assets": [...]}，两种形式：
        - 字符串列表 ["pku/pku-seal.png", ...]：按区域横排等比 contain
        - dict 列表 [{"src","x","y","w","h"}, ...]：绝对坐标（模板原始位置）"""
        assets = region.get("assets") or []
        if not assets:
            return
        if isinstance(assets[0], dict):
            for a in assets:
                ap = SKILL_ROOT / "assets" / a["src"]
                if not ap.exists():
                    warn(f"LOGO 资产不存在: {ap}")
                    continue
                slide.shapes.add_picture(str(ap), Emu(a["x"] * EMU_PER_PX), Emu(a["y"] * EMU_PER_PX),
                                         Emu(a["w"] * EMU_PER_PX), Emu(a["h"] * EMU_PER_PX))
            return
        n = len(assets)
        slot_w = region["w"] / n
        for i, asset in enumerate(assets):
            ap = SKILL_ROOT / "assets" / asset
            if not ap.exists():
                warn(f"LOGO 资产不存在: {ap}")
                continue
            try:
                from PIL import Image as PILImage
                iw, ih = PILImage.open(ap).size
            except Exception:
                iw, ih = 100, 100
            scale = min(slot_w / iw, region["h"] / ih)
            dw, dh = iw * scale, ih * scale
            dx = region["x"] + i * slot_w + (slot_w - dw) / 2
            dy = region["y"] + (region["h"] - dh) / 2
            slide.shapes.add_picture(str(ap), Emu(dx * EMU_PER_PX), Emu(dy * EMU_PER_PX),
                                     Emu(dw * EMU_PER_PX), Emu(dh * EMU_PER_PX))

    def add_smart_run(self, p, seg, size, color_rgb, bold, font_cn, font_en, sup=True):
        """添加 run。sup=True 时 [1] [2,3] 等数字引用标记渲染为上标（baseline +30000，正值=上标）；
        sup=False 时按普通文本渲染（参考文献列表用，[13] 不走上标）。"""
        for sub in re.split(r"(\[\d+(?:[-,，]\d+)*\])", seg):
            if not sub:
                continue
            is_ref = sup and bool(re.fullmatch(r"\[\d+(?:[-,，]\d+)*\]", sub))
            for s2, is_en in split_ascii_runs(sub, font_cn, font_en):
                run = p.add_run()
                run.text = s2
                f = run.font
                f.name = font_en if is_en else font_cn
                f.size = Pt(size)
                f.bold = bold
                f.color.rgb = color_rgb
                if is_ref:
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("baseline", "30000")  # 正值 = 上标
        return

    def render_rules(self, slide, regions):
        """渲染分隔线组：rule 与 rule2（模板内容页为上下两条线）。"""
        for key in ("rule", "rule2"):
            if key in regions:
                r = regions[key]
                self.add_rect(slide, r["x"], r["y"], r["w"], r.get("h", 2), self.pal[r.get("color", "line")])

    def add_image_fit(self, slide, region, img_path):
        x, y, w, h = region["x"], region["y"], region["w"], region["h"]
        try:
            from PIL import Image as PILImage
            iw, ih = PILImage.open(img_path).size
        except Exception:
            iw, ih = 1600, 900
        scale = min(w / iw, h / ih)
        dw, dh = iw * scale, ih * scale
        dx = x + (w - dw) / 2
        dy = y + (h - dh) / 2
        slide.shapes.add_picture(str(img_path), Emu(dx * EMU_PER_PX), Emu(dy * EMU_PER_PX),
                                 Emu(dw * EMU_PER_PX), Emu(dh * EMU_PER_PX))

    def add_table(self, slide, region, headers, rows):
        x, y, w, h = region["x"], region["y"], region["w"], region["h"]
        n_rows, n_cols = len(rows) + 1, len(headers)
        gt = slide.shapes.add_table(n_rows, n_cols, Emu(x * EMU_PER_PX), Emu(y * EMU_PER_PX),
                                    Emu(w * EMU_PER_PX), Emu(h * EMU_PER_PX))
        tbl = gt.table
        tbl.first_row = True
        tbl.horz_banding = True
        body_size = self.size("body", 18)
        for c, head in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = ""
            p0 = cell.text_frame.paragraphs[0]
            self.add_smart_run(p0, str(head), body_size, hex_to_rgb(self.pal["base"]), True,
                               self.font("cn"), self.font("en"))
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(self.pal["accent"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                p0 = cell.text_frame.paragraphs[0]
                self.add_smart_run(p0, str(val), body_size, hex_to_rgb(self.pal["text"]), False,
                                   self.font("cn"), self.font("en"))
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 列宽均分
        for c in range(n_cols):
            tbl.columns[c].width = Emu(int(w / n_cols) * EMU_PER_PX)
        return gt

    def add_stats(self, slide, region, items):
        x, y, w, h = region["x"], region["y"], region["w"], region["h"]
        n = len(items)
        gap = 24
        cw = (w - gap * (n - 1)) / n
        for i, item in enumerate(items):
            cx = x + i * (cw + gap)
            card = self.add_rect(slide, cx, y, cw, h, self.pal["panel"])
            self.add_text(slide, {"x": cx, "y": y + 20, "w": cw, "h": h * 0.45,
                                  "align": "center", "size": "stat-num", "accent": True, "bold": True},
                          str(item.get("num", "")), size_token="stat-num", align="center")
            self.add_text(slide, {"x": cx, "y": y + h * 0.58, "w": cw, "h": h * 0.32,
                                  "align": "center", "size": "stat-label", "muted": True},
                          str(item.get("label", "")), size_token="stat-label", align="center")

    # -- 页面 --

    def render_cover(self, slide, sl, sk):
        bg = sk.get("bg", {})
        self.add_rect(slide, 0, 0, self.canvas["w"], self.canvas["h"], self.pal[bg.get("fill", "base")])
        if bg.get("topAccent"):
            self.add_rect(slide, 0, 0, self.canvas["w"], 8, self.pal["accent"])
        if bg.get("leftAccent"):
            self.add_rect(slide, 0, 0, 10, self.canvas["h"], self.pal["accent"])
        regions = sk["regions"]
        if "ruleTop" in regions:
            r = regions["ruleTop"]
            self.add_rect(slide, r["x"], r["y"], r["w"], r.get("h", 2), self.pal[r.get("color", "line")])
        if "title" in regions:
            self.add_text(slide, regions["title"], sl.get("title") or self.spec["deckTitle"])
        if "subtitle" in regions and (sl.get("subtitle") or self.spec.get("subtitle")):
            self.add_text(slide, regions["subtitle"], sl.get("subtitle") or self.spec.get("subtitle"))
        if "date" in regions and (sl.get("date") or self.spec.get("date")):
            self.add_text(slide, regions["date"], sl.get("date") or self.spec.get("date"))
        if "abstract" in regions and sl.get("abstract"):
            ab = sl["abstract"] if isinstance(sl["abstract"], list) else [sl["abstract"]]
            self.add_text(slide, regions["abstract"], ab)
        if "meta" in regions:
            meta = " · ".join(filter(None, [
                sl.get("author") or self.spec.get("author"),
                self.spec.get("affiliation"),
                sl.get("date") or self.spec.get("date"),
            ]))
            if meta:
                self.add_text(slide, regions["meta"], meta)
        if "ruleBottom" in regions:
            r = regions["ruleBottom"]
            self.add_rect(slide, r["x"], r["y"], r["w"], r.get("h", 2), self.pal[r.get("color", "line")])

    def render_toc(self, slide, sl, sk):
        bg = sk.get("bg", {})
        self.add_rect(slide, 0, 0, self.canvas["w"], self.canvas["h"], self.pal[bg.get("fill", "base")])
        regions = sk["regions"]
        if "tocTitle" in regions:
            self.add_text(slide, regions["tocTitle"], sl.get("title") or "目录")
        elif "header" in regions:
            self.add_text(slide, regions["header"], sl.get("title") or "目录")
        self.render_rules(slide, regions)
        if "logo" in regions:
            self.render_logo(slide, regions["logo"])
        items = sl.get("items") or [f"{s['id']}  {s['title']}" for s in self.spec.get("sections", [])]
        toc_reg = dict(regions["tocItems"])
        lines = [str(it) for it in items]
        self.add_text(slide, toc_reg, lines)
        if "pageNum" in regions:
            self.render_page_num(slide, regions["pageNum"])

    def render_section(self, slide, sl, sk):
        bg = sk.get("bg", {})
        self.add_rect(slide, 0, 0, self.canvas["w"], self.canvas["h"], self.pal[bg.get("fill", "base")])
        if bg.get("leftAccent"):
            self.add_rect(slide, 0, 0, 10, self.canvas["h"], self.pal["accent"])
        regions = sk["regions"]
        if "header" in regions:
            self.add_text(slide, regions["header"], sl.get("title") or "")
        self.render_rules(slide, regions)
        if "logo" in regions:
            self.render_logo(slide, regions["logo"])
        if "sectionNo" in regions:
            sec = self.sections.get(sl.get("section", ""), {})
            no = sl.get("sectionNo") or sec.get("id", "")
            self.add_text(slide, regions["sectionNo"], f"{no:0>2}" if str(no).isdigit() else str(no))
        if "sectionTitle" in regions:
            sec = self.sections.get(sl.get("section", ""), {})
            self.add_text(slide, regions["sectionTitle"], sl.get("title") or sec.get("title", ""))
        if "sectionTagline" in regions and sl.get("tagline"):
            self.add_text(slide, regions["sectionTagline"], sl["tagline"])
        if "pageNum" in regions:
            self.render_page_num(slide, regions["pageNum"])

    def render_content(self, slide, sl, sk, sk_wide):
        layout = sl.get("layout", "bullets")
        has_image = bool(sl.get("image")) and layout in ("image-right", "image-left")
        use_sk = sk if has_image else sk_wide
        if use_sk is None:
            use_sk = sk
        bg = use_sk.get("bg", {})
        self.add_rect(slide, 0, 0, self.canvas["w"], self.canvas["h"], self.pal[bg.get("fill", "base")])
        regions = use_sk["regions"]
        # 页头
        if "header" in regions:
            self.add_text(slide, regions["header"], sl.get("title", ""))
        # 章节归属条
        if "sectionTag" in regions and sl.get("section"):
            sec = self.sections.get(sl["section"], {})
            tag = f"{sec.get('id',''):0>2}·{sec.get('short') or sec.get('title','')}" if str(sec.get('id','')).isdigit() else f"{sec.get('id','')}·{sec.get('short') or sec.get('title','')}"
            self.add_text(slide, regions["sectionTag"], tag)
        # 分隔线（rule + rule2 双线）与右上角 LOGO
        self.render_rules(slide, regions)
        if "logo" in regions:
            self.render_logo(slide, regions["logo"])
        # 主体（先预计算参考文献位置，正文与其重叠时上移+左右拉长）
        body_reg = regions["body"]
        ref_reg, refs_list = None, None
        if "references" in regions and sl.get("references"):
            refs_list = sl["references"] if isinstance(sl["references"], list) else [sl["references"]]
            ref_reg = dict(regions["references"])
            if ref_reg.get("anchor") == "bottom-left":
                ref_size = self.size("ref", 9)
                px = ref_size * 1.333  # 字号像素（9pt ≈ 12px）
                # 按中英混排精确估算行数：ASCII 半角 6px、中文全角 12px
                n_lines = 0
                for rr in refs_list:
                    w_total = sum(px if ord(ch) > 127 else px / 2 for ch in rr)
                    n_lines += max(1, math.ceil(w_total / ref_reg["w"]))
                h_px = int(n_lines * px * 1.15 + 2)  # 行距 1.15 贴近 add_text 实际，底部仅留 2px
                ref_reg["x"] = 0
                ref_reg["y"] = self.canvas["h"] - h_px
                ref_reg["h"] = h_px
            ref_top = ref_reg["y"]
            if body_reg["y"] + body_reg["h"] > ref_top - 8:
                overflow = (body_reg["y"] + body_reg["h"]) - (ref_top - 8)
                min_y = regions.get("rule2", {}).get("y", 80) + 14
                shift = min(overflow, body_reg["y"] - min_y)
                body_reg = dict(body_reg)
                body_reg["y"] -= shift
                body_reg["h"] = (ref_top - 8) - body_reg["y"]
                margin = max(60 - shift // 2, 24)
                body_reg["x"] = margin
                body_reg["w"] = self.canvas["w"] - 2 * margin
                warn(f"正文与参考文献重叠：正文上移 {shift}px、左右拉宽（margin {margin}px）")
        if layout == "two-col":
            cols = sl.get("body", {})
            half = dict(body_reg)
            half["w"] = int(body_reg["w"] / 2) - 20
            self.add_text(slide, half, cols.get("left", []))
            half2 = dict(half)
            half2["x"] += half["w"] + 40
            self.add_text(slide, half2, cols.get("right", []))
        elif layout == "table":
            t = sl.get("body", {})
            self.add_table(slide, body_reg, t.get("headers", []), t.get("rows", []))
        elif layout == "stats":
            self.add_stats(slide, body_reg, sl.get("body", []))
        elif layout == "text":
            body = sl.get("body", [])
            if isinstance(body, str):
                body = [body]
            self.add_text(slide, body_reg, body)
            self.check_overflow(body, body_reg, self.size("body", 18))
        else:  # bullets / image-right / image-left
            body = sl.get("body", [])
            if isinstance(body, str):
                body = [body]
            items = self.parse_bullet_items(body, self.style.get("bullets", {}))
            self.add_rich_paras(slide, body_reg, items)
            self.check_overflow(body, body_reg, self.size("body", 18))
        # 图片（image-left 时左右镜像交换 body/image 位置）
        if has_image and "image" in regions:
            img_reg = regions["image"]
            if layout == "image-left":
                bx, ix = body_reg["x"], img_reg["x"]
                body_reg = dict(body_reg); img_reg = dict(img_reg)
                body_reg["x"], img_reg["x"] = ix, bx
            img = sl["image"]
            img_path = Path(img) if Path(img).is_absolute() else self.spec_dir / img
            if img_path.exists():
                self.add_image_fit(slide, img_reg, str(img_path))
                if sl.get("caption") and "caption" in regions:
                    cap_reg = regions["caption"]
                    if layout == "image-left":
                        cap_reg = dict(cap_reg); cap_reg["x"] = img_reg["x"]
                    self.add_text(slide, cap_reg, sl["caption"])
            else:
                warn(f"图片不存在: {img_path}（已跳过）")
        # 页码 + 页脚（兼容三格/单格两套骨架）
        if "pageNum" in regions:
            self.render_page_num(slide, regions["pageNum"])
        # 参考文献（左下角贴边文本框，[1] 一行一条，正常文本不走上标，位置已在主体前预计算）
        if ref_reg is not None and refs_list is not None:
            self.add_text(slide, ref_reg, refs_list, sup=False)
        if "footerDate" in regions:
            self.add_text(slide, regions["footerDate"], self.spec.get("date", ""))
        if "footerMid" in regions:
            self.add_text(slide, regions["footerMid"], self.spec.get("affiliation", ""))
        if "footer" in regions:
            self.add_text(slide, regions["footer"], f"{self.spec['deckTitle']} · {self.spec.get('author','')}")

    def render_ending(self, slide, sl, sk):
        bg = sk.get("bg", {})
        self.add_rect(slide, 0, 0, self.canvas["w"], self.canvas["h"], self.pal[bg.get("fill", "base")])
        regions = sk["regions"]
        if "ruleTop" in regions:
            r = regions["ruleTop"]
            self.add_rect(slide, r["x"], r["y"], r["w"], r.get("h", 2), self.pal[r.get("color", "line")])
        if "thanks" in regions:
            self.add_text(slide, regions["thanks"], sl.get("title") or "谢谢聆听")
        if "thanksSub" in regions and sl.get("thanks"):
            self.add_text(slide, regions["thanksSub"], sl["thanks"])
        if "contact" in regions and sl.get("contact"):
            self.add_text(slide, regions["contact"], sl["contact"])
        if "ruleBottom" in regions:
            r = regions["ruleBottom"]
            self.add_rect(slide, r["x"], r["y"], r["w"], r.get("h", 2), self.pal[r.get("color", "line")])

    def render_page_num(self, slide, region):
        self.page_no += 1
        self.add_text(slide, region, str(self.page_no))

    def check_overflow(self, text_list, region, size_pt):
        h_pt = estimate_text_height(text_list, size_pt, region["w"])
        region_h_pt = region["h"] * PX_TO_PT
        if h_pt > region_h_pt * 1.02:
            warn(f"正文可能溢出（估算 {h_pt:.0f}pt > 区域 {region_h_pt:.0f}pt）：建议缩短文案/拆页/换布局")

    # -- 主循环 --

    def run(self, spec, out_path):
        sk_map = self.sk["slides"]
        for i, sl in enumerate(spec["slides"]):
            t = sl["type"]
            slide = self.prs.slides.add_slide(self.blank)
            if t == "cover":
                self.render_cover(slide, sl, sk_map.get("cover", {}))
            elif t == "toc":
                self.render_toc(slide, sl, sk_map.get("toc", {}))
            elif t == "section":
                self.render_section(slide, sl, sk_map.get("section", {}))
            elif t == "content":
                self.render_content(slide, sl, sk_map.get("content", {}), sk_map.get("content-wide"))
            elif t == "ending":
                self.render_ending(slide, sl, sk_map.get("ending", {}))
            else:
                warn(f"slide[{i}] type '{t}' 跳过（未知类型）")
                continue
            if sl.get("notes"):
                slide.notes_slide.notes_text_frame.text = sl["notes"]
        out_path = Path(out_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out_path))
        print(f"✓ 已渲染 {len(spec['slides'])} 页 → {out_path}")
        if WARN:
            print(f"⚠️ 共 {len(WARN)} 条警告（不影响产出，QA 阶段重点检查）")


def main():
    ap = argparse.ArgumentParser(description="PPT 专属 skill 渲染引擎")
    ap.add_argument("-s", "--spec", required=True, help="deck-spec.json 路径")
    ap.add_argument("-o", "--output", required=True, help="输出 .pptx 路径")
    ap.add_argument("--style", default=None, help="覆盖风格名（容器风格须指定子风格）")
    args = ap.parse_args()

    spec = load_json(args.spec)
    spec["_spec_path"] = args.spec
    validate_spec(spec)
    style, style_name = resolve_style(spec, args.style)
    skeleton = resolve_skeleton(style)
    print(f"场景={spec.get('scenario')}  风格={style_name}  骨架={skeleton['name']}")
    Renderer(spec, style, skeleton).run(spec, args.output)


if __name__ == "__main__":
    main()
